import os
import re
import shutil
import hashlib
from pathlib import Path
from datetime import datetime

import magic
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
import pytesseract
import pandas as pd

BASE = Path(__file__).resolve().parents[1]
INCOMING = BASE / 'data' / 'incoming'
INDEXED = BASE / 'data' / 'indexed'
INCOMING.mkdir(parents=True, exist_ok=True)
INDEXED.mkdir(parents=True, exist_ok=True)

from app.db import insert_file_meta, get_conn

def sha256_file(path: Path):
    h = hashlib.sha256()
    with path.open('rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            h.update(chunk)
    return h.hexdigest()

def atomic_move(src: Path, dest: Path):
    dest.parent.mkdir(parents=True, exist_ok=True)
    tmp = dest.with_suffix(dest.suffix + '.tmp')
    shutil.copy2(src, tmp)
    os.replace(tmp, dest)

def extract_text_from_pdf(path: Path) -> str:
    try:
        texts = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                texts.append(page.extract_text() or '')
                if i >= 1:
                    break
        return '\n'.join(texts)
    except Exception:
        return ''

def extract_text_from_docx(path: Path) -> str:
    try:
        d = docx.Document(str(path))
        return '\n'.join(p.text for p in d.paragraphs[:200])
    except Exception:
        return ''

def extract_text_from_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        texts = []
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    texts.append(shape.text)
            if i >= 1: break
        return '\n'.join(texts)
    except Exception:
        return ''

def extract_text_from_image(path: Path) -> str:
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img)
    except Exception:
        return ''

def extract_text_from_excel(path: Path) -> str:
    try:
        df = pd.read_excel(path, nrows=20, engine='openpyxl')
        return df.astype(str).head(20).to_csv(index=False)
    except Exception:
        return ''

def generic_text_extract(path: Path, mimetype: str) -> str:
    ext = path.suffix.lower()
    if 'pdf' in mimetype or ext == '.pdf':
        return extract_text_from_pdf(path)
    if ext in ['.docx', '.doc']:
        return extract_text_from_docx(path)
    if ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(path)
    if ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
        return extract_text_from_image(path)
    if ext in ['.xlsx', '.xls', '.csv']:
        return extract_text_from_excel(path)
    try:
        return path.read_text(errors='ignore')[:5000]
    except Exception:
        return ''

def detect_file_type(path: Path, mimetype: str) -> str:
    ext = path.suffix.lower()
    if 'pdf' in mimetype or ext == '.pdf': return 'pdf'
    if ext in ['.xlsx', '.xls', '.csv']: return 'spreadsheet'
    if ext in ['.docx', '.doc']: return 'word'
    if ext in ['.pptx', '.ppt']: return 'presentation'
    if ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']: return 'image'
    return 'other'

DEPT_KEYWORDS = {
    "Computer_Engineering": ["computer engineering", "cse", "comp", "computer"],
    "IT": ["information technology", "it"],
    "Mechanical": ["mechanical", "mech"],
    "Civil": ["civil", "structural"],
    "ENTC": ["electronics", "entc"],
    "AIML": ["artificial intelligence", "machine learning", "ai", "ml"]
}

CATEGORY_KEYWORDS = {
    "Subjects": ["unit", "module", "notes", "question bank", "syllabus", "assignment"],
    "Marks": ["marks", "result", "grade", "score", "marksheet"],
    "Management": ["meeting", "minutes", "budget", "report", "admin"],
    "PersonalData": ["student name", "dob", "admission", "aadhar", "roll no"],
    "Notices": ["notice", "circular", "announcement", "timetable"]
}

SUBJECT_KEYWORDS = {
    "DSA": ["data structure", "dsa", "algorithm"],
    "DBMS": ["dbms", "database", "sql"],
    "OS": ["operating system", "os"],
    "COA": ["computer organization", "coa"],
    "CNS": ["computer networks", "cns", "network"],
    "AI": ["artificial intelligence", "ai"],
    "ML": ["machine learning", "ml"]
}

def _clean_text_for_match(s: str) -> str:
    return (s or "").lower()

def detect_department(text: str, filename: str = "") -> str:
    t = _clean_text_for_match(filename) + " " + _clean_text_for_match(text)
    for dept, keys in DEPT_KEYWORDS.items():
        for k in keys:
            if k in t:
                return dept
    return "Uncategorized"

def detect_year(text: str, filename: str = "") -> str:
    t = _clean_text_for_match(filename) + " " + _clean_text_for_match(text)
    m = re.search(r"(20\d{2})\s*[-/]\s*(\d{2,4})", t)
    if m:
        y1 = m.group(1); y2 = m.group(2)
        if len(y2) == 2:
            return f"{y1}-{y2}"
        return f"{y1}-{y2[-2:]}"
    m2 = re.search(r"(19|20)\d{2}", t)
    if m2:
        return m2.group(0)
    return "UnknownYear"

def detect_category(text: str, filename: str = "") -> str:
    t = _clean_text_for_match(filename) + " " + _clean_text_for_match(text)
    scores = {cat: sum(t.count(k) for k in keys) for cat, keys in CATEGORY_KEYWORDS.items()}
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "Uncategorized"

def detect_subject(text: str, filename: str = "", department: str = "") -> str:
    t = _clean_text_for_match(filename) + " " + _clean_text_for_match(text)
    scores = {subj: sum(t.count(k) for k in keys) for subj, keys in SUBJECT_KEYWORDS.items()}
    best = max(scores, key=scores.get)
    return best if scores[best] > 0 else "General"

def process_file(path: Path):
    path = Path(path)
    if not path.exists():
        return {'ok': False, 'reason': 'not found'}
    mimetype = magic.from_file(str(path), mime=True)
    checksum = sha256_file(path)

    conn = get_conn()
    cur = conn.cursor()
    cur.execute('SELECT id FROM files WHERE checksum = ?', (checksum,))
    if cur.fetchone():
        conn.close()
        return {'ok': True, 'skipped': True, 'reason': 'duplicate'}
    conn.close()

    snippet = generic_text_extract(path, mimetype)
    ftype = detect_file_type(path, mimetype)

    dept = detect_department(snippet, path.name)
    year = detect_year(snippet, path.name)
    category = detect_category(snippet, path.name)

    if category == "Subjects":
        subject = detect_subject(snippet, path.name, dept)
        dest_dir = INDEXED / dept / year / "Subjects" / subject
    else:
        dest_dir = INDEXED / dept / year / category

    dest_dir.mkdir(parents=True, exist_ok=True)
    safe_name = re.sub(r"[^\w\-_.() ]", "_", path.name)
    target = dest_dir / safe_name
    counter = 1
    while target.exists():
        target = dest_dir / f"{target.stem}_{counter}{target.suffix}"
        counter += 1

    atomic_move(path, target)

    file_id = insert_file_meta(
        filename=path.name,
        stored_path=str(target),
        department=dept,
        year=year,
        file_type=ftype,
        snippet=snippet[:4000],
        checksum=checksum,
        created_at=datetime.utcnow().isoformat(),
    )

    return {'ok': True, 'dest': str(target), 'department': dept, 'year': year, 'file_type': ftype, 'file_id': file_id}

def process_all_incoming():
    for f in INCOMING.iterdir():
        if f.is_file():
            print(process_file(f))

if __name__ == '__main__':
    process_all_incoming()
