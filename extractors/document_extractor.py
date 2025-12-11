"""Document text extraction (DOCX, PPTX, XLSX, TXT, etc.)"""
import docx
from pptx import Presentation
from pathlib import Path
import logging
import csv
import json

logger = logging.getLogger(__name__)


class DocumentExtractor:
    """Extract text from document files"""
    
    @staticmethod
    def extract_docx(filepath: Path) -> str:
        """Extract text from DOCX file"""
        try:
            doc = docx.Document(filepath)
            text = '\n'.join([para.text for para in doc.paragraphs])
            return text.strip() if text else ""
        except Exception as e:
            logger.error(f"Error extracting DOCX {filepath}: {e}")
            return ""
    
    @staticmethod
    def extract_pptx(filepath: Path) -> str:
        """Extract text from PPTX (PowerPoint) file"""
        try:
            prs = Presentation(filepath)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"\n=== Slide {slide_num} ===\n"
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            slide_text += row_text + "\n"
                
                text_parts.append(slide_text)
            
            full_text = '\n'.join(text_parts)
            logger.info(f"Extracted {len(prs.slides)} slides from {filepath.name}")
            return full_text.strip() if full_text else f"PowerPoint: {filepath.name}"
            
        except Exception as e:
            logger.error(f"Error extracting PPTX {filepath}: {e}")
            return f"PowerPoint file: {filepath.name}"
    
    @staticmethod
    def extract_pptx_images(filepath: Path, output_dir: Path) -> list:
        """Extract images from PPTX file
        
        Returns: List of extracted image paths with slide info
        """
        try:
            from PIL import Image
            import io
            
            output_dir.mkdir(parents=True, exist_ok=True)
            prs = Presentation(filepath)
            image_paths = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            image_ext = image.ext
                            
                            # Create unique filename
                            image_filename = f"{filepath.stem}_slide{slide_num}_img{shape_idx + 1}.{image_ext}"
                            image_path = output_dir / image_filename
                            
                            # Save image
                            with open(image_path, "wb") as img_file:
                                img_file.write(image_bytes)
                            
                            image_info = {
                                'path': str(image_path),
                                'slide': slide_num,
                                'filename': image_filename
                            }
                            image_paths.append(image_info)
                            logger.info(f"Extracted image from slide {slide_num}: {image_filename}")
                            
                        except Exception as e:
                            logger.warning(f"Could not extract image from slide {slide_num}: {e}")
            
            logger.info(f"Extracted {len(image_paths)} images from {filepath.name}")
            return image_paths
            
        except Exception as e:
            logger.error(f"Error extracting images from PPTX {filepath}: {e}")
            return []
    
    @staticmethod
    def extract_text(filepath: Path) -> str:
        """Extract text from plain text file"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read().strip()
        except Exception as e:
            logger.error(f"Error extracting text {filepath}: {e}")
            return ""
    
    @staticmethod
    def extract_xlsx(filepath: Path) -> str:
        """Extract text from Excel file"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text_parts.append(row_text)
            
            logger.info(f"Extracted {len(wb.sheetnames)} sheets from {filepath.name}")
            return '\n'.join(text_parts).strip()
            
        except ImportError:
            logger.warning("openpyxl not installed, cannot extract Excel files")
            return f"Excel file: {filepath.name}"
        except Exception as e:
            logger.error(f"Error extracting XLSX {filepath}: {e}")
            return f"Excel file: {filepath.name}"
    
    @staticmethod
    def extract_csv(filepath: Path) -> str:
        """Extract text from CSV file"""
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                rows = []
                for i, row in enumerate(reader):
                    if i > 1000:  # Limit to first 1000 rows
                        rows.append(f"... (truncated at 1000 rows)")
                        break
                    rows.append(" | ".join(row))
                return '\n'.join(rows)
        except Exception as e:
            logger.error(f"Error extracting CSV {filepath}: {e}")
            return f"CSV file: {filepath.name}"
    
    @staticmethod
    def extract_json(filepath: Path) -> str:
        """Extract text from JSON file"""
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        except Exception as e:
            logger.error(f"Error extracting JSON {filepath}: {e}")
            return f"JSON file: {filepath.name}"
    
    @staticmethod
    def extract_jupyter(filepath: Path) -> str:
        """Extract text from Jupyter Notebook (.ipynb)"""
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                nb = json.load(f)
                text_parts = []
                
                for i, cell in enumerate(nb.get('cells', [])):
                    cell_type = cell.get('cell_type', '')
                    source = ''.join(cell.get('source', []))
                    
                    if cell_type == 'markdown':
                        text_parts.append(f"\n=== Markdown Cell {i+1} ===\n{source}")
                    elif cell_type == 'code':
                        text_parts.append(f"\n=== Code Cell {i+1} ===\n{source}")
                
                logger.info(f"Extracted {len(nb.get('cells', []))} cells from {filepath.name}")
                return '\n'.join(text_parts).strip()
                
        except Exception as e:
            logger.error(f"Error extracting Jupyter {filepath}: {e}")
            return f"Jupyter Notebook: {filepath.name}"
